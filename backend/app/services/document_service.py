import os
import uuid
from datetime import datetime, timezone
from langchain_community.document_loaders import PyMuPDFLoader, TextLoader, CSVLoader
from typing import Optional, List
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from app.core.config import Config

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".txt", ".csv", ".pptx"}

class DocumentService:
    def __init__(self, upload_dir: Optional[str] = None, indexed_files_log: Optional[str] = None):
        self.upload_dir       = upload_dir       or Config.UPLOAD_DIR
        self.indexed_files_log = indexed_files_log or Config.INDEXED_FILES_LOG

        self.default_text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=Config.CHUNK_SIZE,
            chunk_overlap=Config.CHUNK_OVERLAP,
            separators=["\n\n", "\n", ".", " ", ""]
        )

    def get_domain_from_content(self, text: str) -> str:
        """Heuristic-based domain detection for high-speed categorization."""
        text = text.lower()
        if any(kw in text for kw in ["court", "plaintiff", "defendant", "judgment", "legal", "petitioner", "advocate", "tribunal", "appeal", "verdict"]):
            return "legal"
        if any(kw in text for kw in ["patient", "diagnosis", "symptoms", "medical", "clinical", "hospital", "prescription", "dosage", "therapy", "diagnosis"]):
            return "medical"
        if any(kw in text for kw in ["experience", "education", "skills", "cv", "resume", "work history", "professional summary", "certification"]):
            return "hr"
        if any(kw in text for kw in ["invoice", "bill", "payment", "tax", "financial", "pricing", "revenue", "balance sheet", "audit"]):
            return "financial"
        if any(kw in text for kw in ["methodology", "abstract", "references", "research", "conclusion", "hypothesis", "literature review"]):
            return "academic"
        return "general"

    def _get_text_splitter_for_domain(self, domain: str) -> RecursiveCharacterTextSplitter:
        """Returns a text splitter configured for the given domain."""
        domain_cfg = Config.get_domain_config(domain)
        return RecursiveCharacterTextSplitter(
            chunk_size=domain_cfg["chunk_size"],
            chunk_overlap=domain_cfg["chunk_overlap"],
            separators=["\n\n", "\n", ".", " ", ""]
        )

    def _load_file(self, file_path: str) -> List[Document]:
        """Loads a document from any supported file format."""
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            loader = PyMuPDFLoader(file_path)
            return loader.load()

        elif ext in (".docx", ".doc"):
            import docx2txt
            text = docx2txt.process(file_path)
            return [Document(page_content=text, metadata={"source": file_path, "page": 0})]

        elif ext in (".xlsx", ".xls"):
            return self._load_excel(file_path)

        elif ext == ".txt":
            loader = TextLoader(file_path, encoding="utf-8")
            return loader.load()

        elif ext == ".csv":
            loader = CSVLoader(file_path, encoding="utf-8")
            return loader.load()

        elif ext == ".pptx":
            return self._load_pptx(file_path)

        else:
            raise ValueError(f"Unsupported file extension: {ext}")

    def _load_excel(self, file_path: str) -> List[Document]:
        """Loads Excel files, converting each sheet's rows into document chunks."""
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        documents = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            headers = [str(h) if h is not None else "" for h in rows[0]]
            text_parts = [f"Sheet: {sheet_name}\n"]
            text_parts.append(" | ".join(headers))
            text_parts.append("-" * 40)
            for row in rows[1:]:
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.replace(" | ", "").strip():
                    text_parts.append(row_text)
            documents.append(Document(
                page_content="\n".join(text_parts),
                metadata={"source": file_path, "sheet": sheet_name, "page": 0}
            ))
        wb.close()
        return documents

    def _load_pptx(self, file_path: str) -> List[Document]:
        """Loads PowerPoint files, converting each slide into a document."""
        from pptx import Presentation
        prs = Presentation(file_path)
        documents = []
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            if text_parts:
                documents.append(Document(
                    page_content="\n".join(text_parts),
                    metadata={"source": file_path, "page": i + 1}
                ))
        return documents

    def process_pdf(self, file_path: str):
        """Loads and splits a document into domain-optimized chunks. Supports PDF, DOCX, XLSX, TXT, CSV, PPTX."""
        documents = self._load_file(file_path)

        file_name = os.path.basename(file_path).lower()

        sample_text = "".join([d.page_content for d in documents[:2]])
        detected_domain = self.get_domain_from_content(sample_text)

        text_splitter = self._get_text_splitter_for_domain(detected_domain)
        chunks = text_splitter.split_documents(documents)

        domain_cfg = Config.get_domain_config(detected_domain)
        doc_id = str(uuid.uuid4())
        upload_ts = datetime.now(timezone.utc).isoformat()
        for i, chunk in enumerate(chunks):
            chunk.metadata["chunk_index"] = i
            chunk.metadata["source_file"] = file_name
            chunk.metadata["domain"] = detected_domain
            chunk.metadata["document_id"] = doc_id
            chunk.metadata["upload_timestamp"] = upload_ts

            chunk.metadata["top_k"] = domain_cfg["top_k"]
            chunk.metadata["similarity_threshold"] = domain_cfg["similarity_threshold"]

        print(f"--- DocumentService: {len(chunks)} chunks created for '{file_name}' "
              f"[domain={detected_domain}, chunk_size={domain_cfg['chunk_size']}] ---")
        return chunks

    def get_indexed_files(self):
        """Returns the set of already indexed files for this platform."""
        if not os.path.exists(self.indexed_files_log):
            return set()
        with open(self.indexed_files_log, "r", encoding="utf-8") as f:
            return set(line.strip().lower() for line in f if line.strip())

    def mark_as_indexed(self, filename: str):
        """Logs a filename as indexed for this platform."""
        with open(self.indexed_files_log, "a", encoding="utf-8") as f:
            f.write(f"{filename.lower()}\n")
